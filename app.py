import streamlit as st
import teacherdb as tdb
from pptx import Presentation
import fitz 
import openai
import time
from openai.error import ServiceUnavailableError,RateLimitError
from googleapiclient.discovery import build
import studentdb as sdb
import relation

with open('style.css') as f :
    st.markdown(f'''<style>{f.read()}</style>''',unsafe_allow_html=True)

st.title('INTELLIGRADER')
if 'islogin' not in st.session_state:
    st.session_state['islogin'] = False
    st.session_state['isregister']=False
    st.session_state['Student']=False
    st.session_state['Teacher']=False
    st.session_state['content']=[]
    st.session_state['Summary']=False
    st.session_state['chatbot']=False
    st.session_state['dept']=''
    st.session_state['id']=''
    st.session_state['rollno']=''
    st.session_state['messages'] = []
    st.session_state['show'] = False
    st.session_state['see_yt']=False
def logout():
    st.session_state['islogin'] = False
    st.session_state['isregister']=False
    st.session_state['Student']=False
    st.session_state['Teacher']=False
    st.session_state['content']=[]
    st.session_state['Summary']=False
    st.session_state['chatbot']=False
    st.session_state['dept']=''
    st.session_state['id']=''
    st.session_state['rollno']=''
    st.session_state['messages'] = []
    st.session_state['show'] = False
    st.session_state['see_yt']=False
if not st.session_state['islogin']:
    tab1,tab2=st.tabs(["Student Login","Teacher Login"])
    with tab1:
        st.header('STUDENT AUTHENTICATION 👨🏻‍🎓')
        f={
            'rollno':None,
            'password':None
        }
        with st.form(key="Student_login"):
            f['rollno']=st.text_input("Enter the Roll Number:")
            f['password']=st.text_input("Enter Password:",type="password")
            submit=st.form_submit_button("Login")
        if submit:
            if not all(f.values()):
                st.warning("Fill Completely!")
            else:
                row=sdb.search(f['rollno'],f['password'])
                if row :
                    st.success("Logged in Sucessfully!")
                    st.session_state['islogin']=True
                    st.session_state['isregister']=True
                    st.session_state['rollno']=f['rollno']
                    st.session_state['Student']=True
                    st.rerun()
                else:
                    st.warning("Invalid Credentials")
        if st.button("Sign up",key="studsignup"):
                st.session_state['islogin']=True
                st.session_state['isregister']=False
                st.session_state['Student']=True
                st.session_state['Teacher']=False
                st.rerun()
    with tab2:
        st.header('TEACHER AUTHENTICATION 👨‍🏫')
        f={
            'id':None,
            'password':None
        }
        with st.form(key="Teacher_login"):
            f['id']=st.text_input("Enter the ID:")
            f['password']=st.text_input("Enter Password:",type="password")
            submit=st.form_submit_button("Login")
        if submit:
            if not all(f.values()):
                st.warning("Fill Completely!")
            else:
                row=tdb.search_login(f['id'],f['password'])
                if row:
                    st.success("Logged in Sucessfully!")
                    st.session_state['id']=f['id']
                    st.session_state['dept']=row[3]
                    st.session_state['islogin']=True
                    st.session_state['isregister']=True
                    st.session_state['Teacher']=True
                    st.rerun()
                else:
                    st.warning("Invalid Credentials")
                    st.rerun()

        if st.button("Sign up"):
                st.session_state['islogin']=True
                st.session_state['isregister']=False
                st.session_state['Teacher']=True
                st.session_state['Student']=False
                st.rerun()
elif  not st.session_state['isregister']:    
    if st.session_state['Student']:
        f={
        "rollno":None,
        "password1":None,
        "password2":None,
    }
        with st.form(key='reg-stu-frm') :
            st.subheader('NEW-STUDENT REGISTRATION 💻')
            f['rollno'] = st.text_input('Enter your Roll Number')
            f['password1'] = st.text_input('Enter Password',type='password')
            f['password2'] = st.text_input('Re Enter Password',type='password')
            submit = st.form_submit_button('REGISTER')
        if submit:
            if not all(f.values()):
                st.error("Fill Completely!")
            elif f['password1']!=f['password2']:
                st.error("Password Mismatch!")
            else:
                lst = sdb.search_roll(f['rollno'])
                if not lst :
                    sdb.add_val(f['rollno'],f['password1'])
                    st.success("Registered in Sucessfully!")
                    st.session_state['islogin']=False
                    st.session_state['isregister']=True
                    st.rerun()
                else:
                    st.warning("User Already Exist")
                    time.sleep(1)
                    st.rerun()
        if st.button("Back"):
            st.session_state['islogin']=False
            st.session_state['isregister']=False
            st.session_state['Student']=False
            st.session_state['Teacher']=False
            st.rerun()
    if st.session_state['Teacher']:
        f={
            "id":None,
            "password1":None,
            "password2":None,
            "dept":None
        }
        with st.form(key='tech-reg-frm') :
            st.subheader('NEW-TEACHER REGISTRATION 💻')
            f['id'] = st.text_input('Enter your ID')
            f['dept'] = st.text_input('Enter your Department')
            f['password1'] = st.text_input('Enter Password',type='password')
            f['password2'] = st.text_input('Re Enter Password',type='password')
            submit= st.form_submit_button('REGISTER')
        if submit:
            if not all(f.values()):
                st.error("Fill Completely!")
            elif f['password1']!=f['password2']:
                st.error("Pasword Miss Match!")
            else:
                if not tdb.search(f['id']):
                    tdb.add_val(f['id'],f['password1'],f['dept'])
                    st.success("Registered in Sucessfully!")
                    st.session_state['islogin']=False
                    st.session_state['isregister']=True
                    st.rerun()
                else:
                    st.warning("User Already Exist")
                    time.sleep(1)
                    st.rerun()
        if st.button("Back"):
            st.session_state['islogin']=False
            st.session_state['isregister']=False
            st.session_state['Student']=False
            st.session_state['Teacher']=False
            st.rerun()
else:
    if st.session_state['Student']:

        if not st.session_state['Summary']:
            tab1,tab2= st.tabs(['Summariser','Results'])
            with tab1:

                openai.api_key = st.secrets['openai_apikey1']
                openai.api_base = st.secrets['openai_apibase']

                def extract_all_text_pdf(file):
                    pdf_bytes = file.read()
                    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
                    all_text = []
                    for page in doc:
                        text = page.get_text()
                        if text.strip():
                            all_text.append(text.strip())
                    return "\n".join(all_text)
                cont=[]
                def extract_all_text_ppt(file):
                    prs = Presentation(file)
                    all_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                all_text.append(shape.text.strip())
                    return "\n".join(all_text)

                st.subheader('LECTURE SUMMARISER 🤖')
                ans = st.selectbox('Choose File Format', ['--SELECT--','PPT', 'PDF'])
                file=None
                if ans == 'PPT':
                    file = st.file_uploader('UPLOAD THE PPT FILE', type=['ppt', 'pptx'])
                    if file:
                        cont = extract_all_text_ppt(file)
                        while True:
                            try:
                                response = openai.ChatCompletion.create(
                                    model="llama-3.1-8b-instant",  
                                    messages=[
                                    {"role": "user", "content":'Explain the following content in a clear, structured, and coherent way.Use sub-topics and headings where appropriate.Break complex ideas into simple parts, give analogies if needed, and maintain logical flow.Avoid skipping steps in the explanation.Instructions (bold for emphasis):- Use sub-topics and headings.- Break complex ideas into simple parts.- Provide analogies where helpful.- Maintain logical flow, no skipping steps.- Output should be beginner-friendly but still thorough. Here is the content:'+cont}
                                ]   
                                )
                                break
                            except ServiceUnavailableError:
                                st.warning("Server busy. Retrying in 2 seconds...")
                                time.sleep(2)
                            except RateLimitError as e:
                                st.warning(f"Rate limit reached. Retrying in 12 seconds...")
                                time.sleep(12)
                            except Exception :
                                st.warning('Token Limit Exceeded')
                                break 
                        st.write(response['choices'][0]['message']['content'])
                elif ans == 'PDF':
                    file = st.file_uploader('UPLOAD THE PDF FILE', type=['pdf'])
                    if file:
                        cont = extract_all_text_pdf(file)
                        fg = 1
                        while True:
                            try:
                                response = openai.ChatCompletion.create(
                                    model="llama-3.1-8b-instant",  
                                    messages=[
                                    {"role": "user", "content":'Explain the following content in a clear, structured, and coherent way.Use sub-topics and headings where appropriate.Break complex ideas into simple parts, give analogies if needed, and maintain logical flow.Avoid skipping steps in the explanation.Instructions (bold for emphasis):- Use sub-topics and headings.- Break complex ideas into simple parts.- Provide analogies where helpful.- Maintain logical flow, no skipping steps.- Output should be beginner-friendly but still thorough. Here is the content:'+cont}
                                ]   
                                )
                                break
                            except ServiceUnavailableError:
                                st.warning("Server busy. Retrying in 2 seconds...")
                                time.sleep(2)
                            except RateLimitError as e:
                                st.warning(f"Rate limit reached. Retrying in 12 seconds...")
                                time.sleep(12)
                            except Exception :
                                st.warning('Toekn Limit Exceeded')
                                fg = 0
                                break
                        if fg :
                            st.write(response['choices'][0]['message']['content'])
                next=st.button('CONTENT CHAT')
                if (ans=='--Select--' and next) or ( ans !='--Select--' and next and not file):
                    st.warning("Upload The File!")
                elif next:
                    st.session_state['Summary']=True
                    st.session_state['chatbot']=False
                    st.session_state['content']=cont
                    st.rerun()
                if st.button('Logout',key='summmary'):
                    logout()
                    st.rerun()
            
            with tab2:
                if 'is_done' not in st.session_state:
                    st.session_state['is_done'] = False
                    st.session_state['marks'] = []
                    st.session_state['suggest'] = []
                    st.session_state['dept'] = []
                    st.session_state['selected_sub'] = None
                    st.session_state['indices'] = []
                    st.session_state['selected_paper'] = None

                st.header('RESULTS AND IMPROVEMENTS 📈')

                if not st.session_state['dept']: 
                    resp = relation.search_relation(st.session_state['rollno'])
                    if resp:
                        st.session_state['dept'] = [r[0] for r in resp]
                        st.session_state['marks'] = [r[1] for r in resp]
                        st.session_state['suggest'] = [r[2] for r in resp]

                if not st.session_state['dept']:
                    st.info("No results found for this roll number.")
                else:
                    subjects = list(dict.fromkeys(st.session_state['dept']))

                    with st.form(key='stu-res-frm'):
                        sub = st.selectbox(
                            'Select the Subject',
                            subjects,
                            index=subjects.index(st.session_state['selected_sub'])
                            if st.session_state['selected_sub'] in subjects else 0
                        )
                        is_done = st.form_submit_button('ENTER')

                    if is_done:
                        st.session_state['selected_sub'] = sub
                        st.session_state['indices'] = [
                            i for i, d in enumerate(st.session_state['dept']) if d == sub
                        ]

                if st.session_state['selected_sub']:
                        names = [
                            f"{st.session_state['selected_sub']}-{j+1}"
                            for j in range(len(st.session_state['indices']))
                        ]

                        with st.form(key='subs-form-lst'):
                            paper = st.selectbox(
                                'Select the Paper',
                                names,
                                index=names.index(st.session_state['selected_paper'])
                                if st.session_state['selected_paper'] in names else 0
                            )
                            if st.form_submit_button('SHOW') :
                                st.session_state['show'] = True

                        if st.session_state['show']:
                            st.session_state['selected_paper'] = paper
                            j = int(paper.split('-')[-1]) - 1
                            real_index = st.session_state['indices'][j]

                            st.subheader("Result:")
                            st.write("Marks:", st.session_state['marks'][real_index])
                            st.write("Suggestion:", st.session_state['suggest'][real_index])

                            if st.button("See youtube Explaination!") :
                                    st.session_state['see_yt'] = True 

                            if st.session_state['see_yt'] :        
                                    API_KEY = st.secrets['youtube_apikey']
                                    YOUTUBE_API_SERVICE_NAME = st.secrets['YOUTUBE_API_SERVICE_NAME']
                                    YOUTUBE_API_VERSION = st.secrets['YOUTUBE_API_VERSION']
                                    def youtube_search(query, max_results=1):
                                        youtube = build(YOUTUBE_API_SERVICE_NAME, YOUTUBE_API_VERSION,
                                                        developerKey=API_KEY)

                                        request = youtube.search().list(
                                            q=query,
                                            part="id,snippet",
                                            maxResults=max_results
                                        )
                                        response = request.execute()

                                        results = []
                                        for item in response["items"]:
                                            if item["id"]["kind"] == "youtube#video":
                                                video_id = item["id"]["videoId"]
                                                video_url = f"https://www.youtube.com/watch?v={video_id}"
                                                results.append((item["snippet"]["title"], video_url))
                                        return results
                                    openai.api_key = st.secrets['openai_apikey1']
                                    openai.api_base = st.secrets['openai_apibase']
                                    instructions = """
                                        suggest a youtube search for each in incorrect question 
                                        which should me more realtistic and not fictional ,
                                        most importantly just give the search prompts only not anything else strictly
                                    """

                                    system_prompt = instructions + "\n\n--- PDF Content Start ---\n"+ st.session_state['suggest'][real_index]  + "\n--- PDF Content End ---"
                                    system={'role':'system','content':system_prompt}
                                    user={'role':'user','content':'GIVE ME SEARCH PROMPTS FOR THE RESPECTIVE INCORRECT ANSWERS . ONLY GIVE PROPMTS AND ONE SEARCH PROMPT IN ONE LINE'}
                                    messaged = [system,user]
                                    while True:
                                        try:
                                            response = openai.ChatCompletion.create(
                                                model="llama-3.1-8b-instant",
                                                messages=messaged,
                                            )
                                            break
                                        except ServiceUnavailableError:
                                            st.warning("Server busy. Retrying in 2 seconds...")
                                            time.sleep(2)
                                        except RateLimitError as e:
                                            st.warning(f"Rate limit reached. Retrying in 12 seconds...")
                                            time.sleep(12)                                             
                                    reply = response['choices'][0]['message']['content']
                                    st.session_state['query'] =reply
                                    query=st.session_state['query']
                                    if query.strip():
                                        results = youtube_search(query, max_results=5)
                                        if results:
                                            st.subheader("Search Results:")
                                            for title, url in results:
                                                st.markdown(f"**{title}**  \n[{url}]({url})")
                                                st.video(url)
                                        else:
                                            st.warning("No results found.")
                                    else:
                                        st.error("Please enter a search query.")



                if st.button('Logout',key='evaulation'):
                    logout()
                    st.rerun()
                        
        elif not st.session_state['chatbot']:

            openai.api_key = st.secrets['openai_apikey1']
            openai.api_base = st.secrets['openai_apibase']
            instructions = """
            You are an AI assistant with two modes:
            1. *Factual mode* – When the question is about the provided PDF, answer using its content first.
            2. *Creative fallback mode* – If the answer is not found in the PDF, say: "Not found in document." Then, provide your best answer using your own knowledge or creativity.

            Rules:
            - Always clearly separate the factual part from the creative part.
            - Never fabricate content when claiming it is from the PDF.
            - Keep responses concise unless the user asks for more detail.
            """

            system_prompt = instructions + "\n\n--- PDF Content Start ---\n" + st.session_state['content'] + "\n--- PDF Content End ---"


            if st.session_state['messages']==[]:
                st.session_state['messages'] = [{'role':'system','content':system_prompt}]


            st.subheader('CHAT WITH THE CONTENT 👨🏻‍💻')

            with st.sidebar:
                st.sidebar.title("Model Parameters")
                temperature = st.sidebar.slider('Temperature', min_value=0.0, max_value=2.0, value=0.7, step=0.1)
                max_tokens = st.sidebar.slider('Token Limit', min_value=1, max_value=4096, value=3000, step=64)
                if st.sidebar.button("Change File"):
                    st.session_state['Summary']=False
                    st.session_state['messages']=[]
                    st.session_state['content']=[]
                    st.rerun()
                if st.sidebar.button('Logout'):
                    logout()
                    st.rerun()

            for message in st.session_state['messages'][1:]:
                with st.chat_message(message['role']):
                    st.markdown(message['content'])

            prompt = st.chat_input('Ask Your Doubt')
            if prompt:
                st.session_state['messages'].append({"role": "user", "content": prompt})

                while True:
                    try:
                        response = openai.ChatCompletion.create(
                            model="llama-3.1-8b-instant",
                            messages=st.session_state['messages'],
                            temperature=temperature,
                            max_tokens=max_tokens
                        )
                        break
                    except ServiceUnavailableError:
                        st.warning("Server busy. Retrying in 2 seconds...")
                        time.sleep(2)
                    except RateLimitError as e:
                        st.warning(f"Rate limit reached. Retrying in 12 seconds...")
                        time.sleep(12)
                        
                reply = response['choices'][0]['message']['content']
                st.session_state['messages'].append({"role": "assistant", "content": reply})
                for message in st.session_state['messages'][-2:]:
                    with st.chat_message(message['role']):
                        st.markdown(message['content'])
                st.rerun()

                
    elif st.session_state['Teacher']:
        def extract_all_text_pdf(file):
            pdf_bytes = file.read()
            doc = fitz.open(stream=pdf_bytes, filetype="pdf")
            all_text = []
            for page in doc:
                text = page.get_text()
                if text.strip():
                    all_text.append(text.strip())
            return "\n".join(all_text)

        if 'submit1' not in st.session_state :
            st.session_state['submit1'] = False 

        if 'submit2' not in st.session_state :
            st.session_state['submit2'] = False 

        st.header('STUDENT EVALUATOR 👨‍🏫')

        rolls = sdb.fetch_rolls()
        dept = st.session_state['dept']
        print(dept)
        roll= None
        with st.form(key='ev-frm'):
            roll = st.selectbox('Select the Student-Roll',rolls)
            ans = st.selectbox(
                'Choose mode of Evaluation',
                ['Student-Paper & Answer-Sheet', 'Student-Paper']
            )
            sub=st.form_submit_button('SUBMIT')
            if sub:
                st.session_state['submit1']=True

        if st.session_state['submit1'] :
            if ans == "Student-Paper & Answer-Sheet":
                with st.form(key='q-form'):
                    st.subheader('Upload Student Paper and Answer Sheet 🧾')
                    script = st.file_uploader('Upload Student Paper', type=['pdf'])
                    key = st.file_uploader('Upload Answer Sheet', type=['pdf'])
                    submit2 = st.form_submit_button('UPLOAD')
                if submit2:
                    st.session_state['submit2']=True
            
                if st.session_state['submit2'] :
                    if not key or not script :
                        st.warning('Please Upload Both Files')
                        st.session_state['submit2'] = False 
                    else :
                        paper = extract_all_text_pdf(script)
                        answer = extract_all_text_pdf(key)
                        system_prompt = f"""MAIN NOTE -> THE TEXT IS PARASED FROM IMAGE SO IT WILL BE DISTORTED USE YOUR WHOLE TALENT TO UNDERSTAND IT
                                            You are an evaluation assistant for student papers and experienced about 20 years in {dept}.
                                            if the subject is maths and if person solved wrong then give how to solve it correctly also
                                            Assgin marks correctly dont give randomly take a serious note on it
                                            When given the student's answers and the correct answer sheet, your job is to evaluate and return the result in the exact format below:
                                            Format Rules (must follow exactly):
                                            1. First line: Write the total marks scored by the student out of 100 if not mentioned , else what is mentioned in the evaluvation script is considered . in the format:
                                            Marks: <score>/Total Marks
                                            2. Second and third lines: Leave completely empty.
                                            3. From the fourth line onward: For each question, provide a short, constructive suggestion describing where the student should improve and any mistakes they made. Mention the question number first.
                                            - Example: Q1: You missed mentioning the key formula; revise the derivation steps.
                                            - Keep feedback clear and specific to that question.
                                            4. Do not include any explanations about your reasoning, only return the evaluation in the required format.
                                            5. Use concise and professional language.
                                            6. If the student did not answer a question, still include a suggestion for that question (e.g., "Q5: No answer provided, revise topic XYZ").Always follow this structure exactly."""
                        messaged = [{'role':'system','content':system_prompt},
                                    {'role':'user','content':f"Student Script:\n{paper}\n\nAnswer Sheet:\n{answer}"}]
                        openai.api_key = st.secrets['openai_apikey2']
                        openai.api_base = st.secrets['openai_apibase']
                        response = openai.ChatCompletion.create(
                        model="llama-3.1-8b-instant",  
                        messages = messaged
                        )      
                        ans = response['choices'][0]['message']['content']
                        st.text(ans)
                        marks = ans.split('\n')[0]
                        sug = ans.split('\n')[1:]                      
                        l = ''
                        for s in sug :
                            l += s
                        relation.add_relation(st.session_state['id'],st.session_state['dept'],roll,marks,l)


            elif ans == "Student-Paper":
                with st.form(key='q,a-form'):
                    st.subheader('Upload Student Paper')
                    script = st.file_uploader('Upload Student Paper', type=['pdf'])
                    submit2 = st.form_submit_button('UPLOAD')

                if submit2:
                    if not script:
                        st.warning('⚠ Please upload the student paper')
                    else:
                        paper = extract_all_text_pdf(script)
                        system_prompt = f"""MAIN NOTE -> THE TEXT IS PARASED FROM IMAGE SO IT WILL BE DISTORTED USE YOUR WHOLE TALENT TO UNDERSTAND IT
                                            You are an evaluation assistant for student papers and experienced about 20 years in {dept}.
                                            Assgin marks correctly dont give randomly take a serious note on it
                                            if the subject is maths and if person solved wrong then give how to solve it correctly also see every single and single and single step clearly all numbers 
                                            When given the student's answers,your job is to evaluate and return the result in the exact format below:
                                            Format Rules (must follow exactly):
                                            1. First line: Write the total marks scored by the student out of 100 if not mentioned , else what is mentioned in the evaluvation script is considered . in the format:
                                            Marks: <score>/Total Marks . IMPORTANT FOLLOW THIS
                                            2. Second and third lines: Leave completely empty.
                                            3. From the fourth line onward: For each question, provide a short, constructive suggestion describing where the student should improve and any mistakes they made. Mention the question number first.
                                            - Example: Q1: You missed mentioning the key formula; revise the derivation steps.
                                            - Keep feedback clear and specific to that question.
                                            4. Do not include any explanations about your reasoning, only return the evaluation in the required format.
                                            5. Use concise and professional language.
                                            6. If the student did not answer a question, still include a suggestion for that question (e.g., "Q5: No answer provided, revise topic XYZ").Always follow this structure exactly."""
                        messaged = [{'role':'system','content':system_prompt},
                                    {'role':'user','content':f"Student Script:\n\n{paper}\n"}]
                        openai.api_key = st.secrets['openai_apikey2']
                        openai.api_base = st.secrets['openai_apibase']
                        response = openai.ChatCompletion.create(
                        model="llama-3.1-8b-instant",  
                        messages = messaged
                        )      
                        ans = response['choices'][0]['message']['content']
                        st.text(ans)
                        marks = ans.split('\n')[0]
                        sug = ans.split('\n')[1:]
                        l = ''
                        for s in sug :
                            l += s
                        relation.add_relation(st.session_state['id'],st.session_state['dept'],roll,marks,l)
        if st.button('Logout'):
            logout() 
            st.rerun()
